from loguru import logger
import os
from io import BufferedReader
from typing import Optional
from fastapi import UploadFile
import mimetypes
from PyPDF2 import PdfReader
import docx2txt
import csv
import pptx

from models.models import Document, DocumentMetadata, Source
from services.extract_metadata import extract_metadata_from_document


async def get_document_from_file(file: UploadFile) -> Document:
    extracted_text = await extract_text_from_form_file(file)
    extracted_metadata = extract_metadata_from_document(extracted_text)
    logger.info(f"Extracted metadata: {extracted_metadata}")

    metadata = DocumentMetadata(source    = Source.file, 
                                source_id = file.filename, 
                                **extracted_metadata)
    logger.info(metadata)
    doc = Document(text=extracted_text, metadata=metadata)
    
    return doc


def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None) -> str:
    """Return the text content of a file given its filepath."""

    if mimetype is None:
        # Get the mimetype of the file based on its extension
        mimetype, _ = mimetypes.guess_type(filepath)

    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise Exception("Unsupported file type")

    # Open the file in binary mode
    file = open(filepath, "rb")
    extracted_text = extract_text_from_file(file, mimetype)

    return extracted_text


excel_mimetypes = ["application/vnd.ms-excel",                                           # Excel 97-2003 Workbook (.xls)
                   "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # Excel Workbook (.xlsx)
                   "application/vnd.ms-excel.sheet.macroEnabled.12",                     # Excel Macro-Enabled Workbook (.xlsm)
                   "application/vnd.ms-excel.sheet.binary.macroEnabled.12",              # Excel Binary Workbook (.xlsb) 
                   "application/vnd.ms-excel.template.macroEnabled.12",                  # Excel Template (.xlt)
                   "application/vnd.ms-excel.template.macroEnabled.12",                  # Excel Macro-Enabled Template (.xltm)
                   "application/vnd.ms-excel.addin.macroEnabled.12"]                     # Excel Add-In (.xlam)



def extract_text_from_file(file: BufferedReader, mimetype: str) -> str:
    if mimetype == "application/pdf":
        # Extract text from pdf using PyPDF2
        reader = PdfReader(file)
        extracted_text = " ".join([page.extract_text() for page in reader.pages])
    elif mimetype == "text/plain" or mimetype == "text/markdown":
        # Read text from plain text file
        extracted_text = file.read().decode("utf-8")
    elif mimetype in excel_mimetypes:
        # really want the original filename here not a fp since not clear if the file suffix is in play
        # libreoffice --headless --convert-to csv --outdir /home/user/documents /home/user/documents/input_file.xlsm
        raise ValueError("Excel files are not yet supported.")
    elif (mimetype == "application/msword"):
        input_file = "/tmp/tmp.doc"
        open(input_file, 'wb').write(file.read())
        import subprocess       
        output_folder = "/tmp/"
        output_format = "docx:\"Office Open XML Text\""
        command = f"libreoffice --headless --convert-to {output_format} --outdir {output_folder} {input_file}"
        try:
            subprocess.run(command, shell=True, check=True)
        except:
            raise ValueError("Unable to convert doc to docx.")      
        extracted_text = docx2txt.process(f"{output_folder}/tmp.docx")  
        os.unlink(f"{output_folder}/tmp.docx")
        os.unlink(input_file)
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
    elif mimetype == "text/csv":
        # Extract text from csv using csv module
        extracted_text = ""
        decoded_buffer = (line.decode("utf-8") for line in file)
        reader = csv.reader(decoded_buffer)
        for row in reader:
            extracted_text += " ".join(row) + "\n"
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    else:
        # Unsupported file type
        file.close()
        raise ValueError("Unsupported file type: {}".format(mimetype))

    file.close()
    return extracted_text


# Extract text from a file based on its mimetype
async def extract_text_from_form_file(file: UploadFile):
    """Return the text content of a file."""
    # get the file body from the upload file object
    mimetype = file.content_type
    logger.info(f"mimetype: {mimetype}")

    file_stream = await file.read()

    temp_file_path = "/tmp/temp_file"

    # write the file to a temporary location
    with open(temp_file_path, "wb") as f:
        f.write(file_stream)

    try:
        extracted_text = extract_text_from_filepath(temp_file_path, mimetype)
    except Exception as e:
        logger.error(f"Error: {e}")
        os.remove(temp_file_path)
        raise e

    # remove file from temp location
    os.remove(temp_file_path)

    return extracted_text
